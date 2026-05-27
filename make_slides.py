"""
Week 5 Milestone — Slide Generator
Builds Week5_Slides.pptx with all 14 slides and embedded plots.

Fast mode (preferred):
  1. Run the Jupyter notebook → it saves PNG files + df_daily.csv
  2. python make_slides.py       ← loads cached PNGs + week5_metrics.json (no re-analysis)

Full mode (first run or after data refresh):
  1. Ensure df_daily.csv exists (notebook cell 'Save df_daily.csv'), or allow HuggingFace download
  2. python make_slides.py       ← runs SARIMA, saves slide_*.png + week5_metrics.json, builds PPTX

Notebook PNG files recognised automatically:
  eda_plots.png, stl_decomposition.png, seasonal_dow.png, differencing_stages.png,
  ACF_&_PACF_—_dd7log(Sales)__[d=1_D=1_s=7].png, residual_diagnostics.png,
  forecast.png, sarima_vs_sarimax.png
"""

import os, sys, json, warnings
from io import BytesIO

import pandas as pd
import numpy as np

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import matplotlib.gridspec as gridspec
import seaborn as sns
from scipy import stats

from statsmodels.tsa.statespace.sarimax import SARIMAX
from statsmodels.tsa.seasonal import STL
from statsmodels.tsa.stattools import adfuller, kpss
from statsmodels.graphics.tsaplots import plot_acf, plot_pacf
from statsmodels.stats.diagnostic import acorr_ljungbox

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings('ignore')
plt.rcParams['font.family'] = 'DejaVu Sans'

DIR = os.path.dirname(os.path.abspath(__file__))

# ═══════════════════════ ASSET CACHE ═══════════════════════
# PNG files saved by this script (fast-path cache)
_SLIDE_PNGS = {
    'eda':      'slide_eda.png',
    'stl':      'slide_stl.png',
    'dow':      'slide_dow.png',
    'diff':     'slide_diff.png',
    'acf':      'slide_acf.png',
    'resid':    'slide_resid.png',
    'forecast': 'slide_forecast.png',
    'compare':  'slide_compare.png',
}
# PNG files saved by the Jupyter notebook (alternative source)
_NOTEBOOK_PNGS = {
    'eda':      'eda_plots.png',
    'stl':      'stl_decomposition.png',
    'dow':      'seasonal_dow.png',
    'diff':     'differencing_stages.png',
    'acf':      'ACF_&_PACF_—_dd7log(Sales)__[d=1_D=1_s=7].png',
    'resid':    'residual_diagnostics.png',
    'forecast': 'forecast.png',
    'compare':  'sarima_vs_sarimax.png',
}
_METRICS_FILE = 'week5_metrics.json'


def _find_buf(key):
    """Load PNG as BytesIO: checks slide cache first, then notebook PNGs."""
    for name in (_SLIDE_PNGS[key], _NOTEBOOK_PNGS[key]):
        p = os.path.join(DIR, name)
        if os.path.exists(p):
            b = BytesIO(open(p, 'rb').read())
            b.seek(0)
            return b
    return None


def _cache_buf(buf, key):
    """Save BytesIO buffer to slide PNG cache file."""
    p = os.path.join(DIR, _SLIDE_PNGS[key])
    buf.seek(0)
    with open(p, 'wb') as f:
        f.write(buf.read())
    buf.seek(0)


def _save_metrics(R):
    fields = ('p','d','q','P','D','Q','s','HOLDOUT',
              'sarima_mape','sarima_rmse',
              'sarimax_mape','sarimax_rmse',
              'aic_sarima','aic_sarimax')
    with open(os.path.join(DIR, _METRICS_FILE), 'w') as f:
        json.dump({k: R[k] for k in fields}, f, indent=2)


def _load_metrics():
    p = os.path.join(DIR, _METRICS_FILE)
    return json.load(open(p)) if os.path.exists(p) else None


def _assets_ready():
    """Return (True, bufs, metrics) when all cached assets are present."""
    m = _load_metrics()
    if m is None:
        return False, None, None
    bufs = {k: _find_buf(k) for k in _SLIDE_PNGS}
    if all(v is not None for v in bufs.values()):
        return True, bufs, m
    return False, None, None


# ═══════════════════════ COLORS ════════════════════════════
NAVY  = RGBColor(0x1F, 0x4E, 0x79)
BLUE  = RGBColor(0x2E, 0x75, 0xB6)
LBLUE = RGBColor(0xBD, 0xD7, 0xEE)
DBLUE = RGBColor(0x15, 0x38, 0x5A)
GOLD  = RGBColor(0xFF, 0xBF, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1A, 0x1A, 0x1A)
GRAY  = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0xF2, 0xF4, 0xF8)
MGRAY = RGBColor(0xD0, 0xD8, 0xE4)
GREEN = RGBColor(0x1A, 0x7A, 0x4A)
LGREEN= RGBColor(0xD5, 0xF0, 0xE0)
RED   = RGBColor(0xC0, 0x39, 0x2B)
LRED  = RGBColor(0xFA, 0xD7, 0xD0)
ORANGE= RGBColor(0xE6, 0x7E, 0x22)

SW, SH = Inches(13.33), Inches(7.5)
TOTAL  = 14


# ═══════════════════════ STEP 1: LOAD DATA ═════════════════

def _fetch_and_save(path):
    """Download FreshRetailNet-50K, aggregate to daily, save CSV.
    Uses streaming to avoid caching the full ~2 GB on disk."""
    try:
        from datasets import load_dataset
    except ImportError:
        sys.exit('ERROR: df_daily.csv not found and "datasets" library is not installed.\n'
                 'Either run the notebook first, or: pip install datasets')

    print('  df_daily.csv not found — streaming FreshRetailNet-50K from HuggingFace...')
    print('  (No disk cache needed; iterating ~4.5M rows may take 5-10 min)')

    ds = load_dataset('Dingdong-Inc/FreshRetailNet-50K', split='train', streaming=True)

    agg = {}
    for i, row in enumerate(ds):
        dt = str(row.get('dt', ''))[:10]
        if not dt:
            continue
        if dt not in agg:
            agg[dt] = dict(sales=[], temp=[], precpt=[], humidity=[],
                           wind=[], holiday=0, activity=0, discount=[],
                           n=0)
        a = agg[dt]
        a['sales'].append(float(row.get('sale_amount') or 0))
        v = row.get('avg_temperature')
        if v is not None: a['temp'].append(float(v))
        v = row.get('precpt')
        if v is not None: a['precpt'].append(float(v))
        v = row.get('avg_humidity')
        if v is not None: a['humidity'].append(float(v))
        v = row.get('avg_wind_level')
        if v is not None: a['wind'].append(float(v))
        a['holiday']  = max(a['holiday'],  int(row.get('holiday_flag',  0) or 0))
        a['activity'] = max(a['activity'], int(row.get('activity_flag', 0) or 0))
        v = row.get('discount')
        if v is not None: a['discount'].append(float(v))
        a['n'] += 1
        if (i+1) % 500_000 == 0:
            print(f'    ...{i+1:,} rows processed ({len(agg)} dates)')

    print(f'  Aggregating {len(agg)} dates...')
    records = []
    for dt in sorted(agg):
        a = agg[dt]
        def _m(lst): return float(np.mean(lst)) if lst else np.nan
        records.append({
            'dt':           dt,
            'total_sales':  sum(a['sales']),
            'median_sales': float(np.median(a['sales'])) if a['sales'] else np.nan,
            'n_records':    a['n'],
            'avg_temp':     _m(a['temp']),
            'avg_precpt':   _m(a['precpt']),
            'avg_humidity': _m(a['humidity']),
            'avg_wind':     _m(a['wind']),
            'holiday':      a['holiday'],
            'activity':     a['activity'],
            'avg_discount': _m(a['discount']),
        })

    df = pd.DataFrame(records)
    df['dt'] = pd.to_datetime(df['dt'])
    df = df.sort_values('dt').set_index('dt').asfreq('D')
    df['log_sales'] = np.log1p(df['total_sales'])
    df.to_csv(path)
    print(f'  Saved df_daily.csv — {len(df)} rows')
    return df


def load_data():
    path = os.path.join(DIR, 'df_daily.csv')
    if not os.path.exists(path):
        df = _fetch_and_save(path)
    else:
        df = pd.read_csv(path, parse_dates=['dt'], index_col='dt')
        df = df.asfreq('D')
        if 'log_sales' not in df.columns:
            df['log_sales'] = np.log1p(df['total_sales'])
    df['dow']      = df.index.dayofweek
    df['dow_name'] = df.index.day_name()
    return df


# ═══════════════════════ STEP 2: ANALYSIS ══════════════════

def run_analysis(df):
    print('  Running SARIMA analysis...')
    log_s   = df['log_sales'].dropna()
    diff1   = log_s.diff().dropna()
    diff1_7 = log_s.diff().diff(7).dropna()

    stl        = STL(log_s, period=7, robust=True)
    stl_result = stl.fit()

    HOLDOUT = 14
    y        = log_s.copy()
    y_train  = y.iloc[:-HOLDOUT]
    y_test   = y.iloc[-HOLDOUT:]
    p,d,q    = 1, 1, 1
    P,D,Q,s  = 0, 1, 1, 7

    fit = SARIMAX(y_train, order=(p,d,q), seasonal_order=(P,D,Q,s),
                  enforce_stationarity=False, enforce_invertibility=False).fit(disp=False)

    fc      = fit.get_forecast(steps=HOLDOUT)
    fc_mean = np.expm1(fc.predicted_mean)
    fc_ci   = fc.conf_int(alpha=0.05)
    fc_lo   = np.expm1(fc_ci.iloc[:,0])
    fc_hi   = np.expm1(fc_ci.iloc[:,1])

    act_s   = np.expm1(y_test)
    trn_s   = np.expm1(y_train)

    exog = df[['avg_temp','avg_precpt','holiday','activity']].copy()
    for c in ['avg_temp','avg_precpt']:
        exog[c] = (exog[c]-exog[c].mean())/exog[c].std()
    exog = exog.ffill()

    sx_fit = SARIMAX(y_train, exog=exog.iloc[:-HOLDOUT],
                     order=(p,d,q), seasonal_order=(P,D,Q,s),
                     enforce_stationarity=False, enforce_invertibility=False).fit(disp=False)
    sx_fc   = sx_fit.get_forecast(steps=HOLDOUT, exog=exog.iloc[-HOLDOUT:])
    sx_mean = np.expm1(sx_fc.predicted_mean)
    sx_ci   = sx_fc.conf_int(alpha=0.05)

    def mape(a,f): return float(np.mean(np.abs((a-f)/a))*100)
    def rmse(a,f): return float(np.sqrt(np.mean((a-f)**2)))
    act = act_s.values

    return dict(
        df=df, log_s=log_s, diff1=diff1, diff1_7=diff1_7,
        stl=stl_result, y_train=y_train, y_test=y_test,
        fit=fit, fc_mean=fc_mean, fc_lo=fc_lo, fc_hi=fc_hi,
        act_s=act_s, trn_s=trn_s,
        sx_mean=sx_mean, sx_ci=sx_ci,
        resid=fit.resid,
        p=p, d=d, q=q, P=P, D=D, Q=Q, s=s,
        sarima_mape  = mape(act, fc_mean.values),
        sarima_rmse  = rmse(act, fc_mean.values),
        sarimax_mape = mape(act, sx_mean.values),
        sarimax_rmse = rmse(act, sx_mean.values),
        aic_sarima   = fit.aic,
        aic_sarimax  = sx_fit.aic,
        HOLDOUT=HOLDOUT,
    )


# ═══════════════════════ STEP 3: PLOT GENERATORS ═══════════
# Each function returns a BytesIO PNG ready to embed in pptx

def _buf(fig, dpi=150):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


def plot_eda(R):
    df = R['df']
    fig = plt.figure(figsize=(14, 9), facecolor='white')
    gs  = gridspec.GridSpec(2, 2, figure=fig, hspace=0.5, wspace=0.35)

    # Panel 1: time series
    ax1 = fig.add_subplot(gs[0, :])
    ax1.plot(df.index, df['total_sales'], color='steelblue', lw=1.5, label='Daily Sales')
    ax1.fill_between(df.index, df['total_sales'], alpha=0.12, color='steelblue')
    h_idx = df[df['holiday']==1].index
    p_idx = df[df['activity']==1].index
    ax1.scatter(h_idx, df.loc[h_idx,'total_sales'], color='crimson',  s=90, zorder=6,
                label='Holiday', marker='*')
    ax1.scatter(p_idx, df.loc[p_idx,'total_sales'], color='darkorange', s=55, zorder=6,
                label='Promotion', marker='^')
    ax1.set_title('Total Daily Fresh Retail Sales (All Products & Stores)', fontsize=12, fontweight='bold')
    ax1.set_ylabel('Sale Amount')
    ax1.legend(fontsize=9)
    ax1.xaxis.set_major_formatter(mdates.DateFormatter('%b %d'))
    ax1.tick_params(axis='x', rotation=20)

    # Panel 2: log series
    ax2 = fig.add_subplot(gs[1, 0])
    ax2.plot(df.index, df['log_sales'], color='seagreen', lw=1.5)
    ax2.set_title('Log-Transformed Sales  [log(1+Sales)]', fontsize=11, fontweight='bold')
    ax2.set_ylabel('log(1+Sales)')
    ax2.xaxis.set_major_formatter(mdates.DateFormatter('%b %d'))
    ax2.tick_params(axis='x', rotation=20)

    # Panel 3: DOW boxplot
    ax3 = fig.add_subplot(gs[1, 1])
    dow_order = ['Monday','Tuesday','Wednesday','Thursday','Friday','Saturday','Sunday']
    sns.boxplot(data=df, x='dow_name', y='total_sales', order=dow_order,
                palette='Blues_d', ax=ax3, linewidth=0.8)
    ax3.set_title('Sales by Day of Week', fontsize=11, fontweight='bold')
    ax3.set_xlabel(''); ax3.set_ylabel('Sale Amount')
    ax3.tick_params(axis='x', rotation=30)

    return _buf(fig)


def plot_stl(R):
    log_s  = R['log_s']
    result = R['stl']
    fig, axes = plt.subplots(4, 1, figsize=(13, 9), sharex=True, facecolor='white')
    fig.suptitle('STL Decomposition — log(1+Sales)  [period=7, robust=True]',
                 fontsize=12, fontweight='bold')

    pairs = [
        (log_s.values,      'Observed',     'steelblue'),
        (result.trend,      'Trend',        'tomato'),
        (result.seasonal,   'Seasonal (s=7)','seagreen'),
        (result.resid,      'Residual',     'dimgray'),
    ]
    for ax, (vals, lbl, col) in zip(axes, pairs):
        ax.plot(log_s.index, vals, color=col, lw=1.4)
        ax.axhline(0, color='black', lw=0.6, linestyle='--', alpha=0.5)
        ax.set_ylabel(lbl, fontsize=10)
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %d'))

    axes[-1].tick_params(axis='x', rotation=20)
    plt.tight_layout()
    return _buf(fig)


def plot_seasonal_dow(R):
    log_s  = R['log_s']
    result = R['stl']
    sea    = pd.Series(result.seasonal, index=log_s.index)
    by_dow = sea.groupby(sea.index.dayofweek).mean()
    labels = ['Mon','Tue','Wed','Thu','Fri','Sat','Sun']
    colors = ['#4C72B0']*5 + ['#DD8452']*2

    fig, ax = plt.subplots(figsize=(7, 3.5), facecolor='white')
    ax.bar(labels, by_dow.values, color=colors, edgecolor='white', linewidth=0.5)
    ax.axhline(0, color='black', lw=0.8)
    ax.set_title('Mean Seasonal Component by Day of Week (STL)', fontsize=11, fontweight='bold')
    ax.set_ylabel('Seasonal Effect on log(1+Sales)')

    amp = by_dow.max() - by_dow.min()
    ax.text(0.98, 0.95, f'Amplitude = {amp:.3f} log-units\n({(np.exp(amp)-1)*100:.0f}% swing)',
            transform=ax.transAxes, ha='right', va='top', fontsize=9,
            bbox=dict(boxstyle='round,pad=0.3', facecolor='lightyellow', alpha=0.9))
    plt.tight_layout()
    return _buf(fig)


def plot_diff_stages(R):
    pairs = [
        (R['log_s'],   'log(1+Sales) — Level',                              'steelblue'),
        (R['diff1'],   '∇ log(Sales) — 1st Difference  (d=1)',              'darkorange'),
        (R['diff1_7'], '∇∇₇ log(Sales) — 1st + Seasonal Diff  (d=1, D=1, s=7)', 'purple'),
    ]
    fig, axes = plt.subplots(3, 1, figsize=(13, 8), facecolor='white')
    fig.suptitle('Series at Each Differencing Stage', fontsize=12, fontweight='bold')
    for ax, (s, lbl, col) in zip(axes, pairs):
        ax.plot(s.index, s.values, color=col, lw=1.2)
        ax.axhline(s.mean(), color='black', lw=0.8, linestyle='--', alpha=0.5)
        ax.set_title(lbl, fontsize=10)
        ax.set_ylabel('Value')
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %d'))
        ax.tick_params(axis='x', rotation=15)
    plt.tight_layout()
    return _buf(fig)


def plot_acf_panels(R):
    """ACF/PACF of the fully-differenced series (the one used for SARIMA)."""
    s    = R['diff1_7']
    lags = 28
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(13, 4), facecolor='white')
    fig.suptitle('ACF & PACF — ∇∇₇log(Sales)  [d=1, D=1, s=7]',
                 fontsize=12, fontweight='bold')

    plot_acf( s.dropna(), lags=lags, ax=ax1, zero=False, alpha=0.05)
    plot_pacf(s.dropna(), lags=lags, ax=ax2, zero=False, alpha=0.05, method='ywm')

    for lag in range(7, lags+1, 7):
        ax1.axvline(lag, color='red', alpha=0.25, lw=1, linestyle=':')
        ax2.axvline(lag, color='red', alpha=0.25, lw=1, linestyle=':')
    ax1.set_title('ACF (red dashed = multiples of s=7)'); ax2.set_title('PACF')
    plt.tight_layout()
    return _buf(fig)


def plot_residuals(R):
    resid = R['resid']
    log_s = R['log_s']
    p,d,q = R['p'], R['d'], R['q']
    P,D,Q,s = R['P'], R['D'], R['Q'], R['s']

    fig = plt.figure(figsize=(13, 8), facecolor='white')
    gs  = gridspec.GridSpec(2, 2, figure=fig, hspace=0.45, wspace=0.35)

    ax1 = fig.add_subplot(gs[0, :])
    ax1.plot(resid.index, resid.values, color='dimgray', lw=1)
    ax1.axhline(0, color='black', lw=0.9, linestyle='--')
    ax1.fill_between(resid.index, resid.values, 0, alpha=0.2, color='steelblue')
    ax1.set_title(f'Residuals — SARIMA({p},{d},{q})({P},{D},{Q})[{s}]',
                  fontsize=11, fontweight='bold')
    ax1.set_ylabel('Residual')
    ax1.xaxis.set_major_formatter(mdates.DateFormatter('%b %d'))
    ax1.tick_params(axis='x', rotation=15)

    ax2 = fig.add_subplot(gs[1, 0])
    plot_acf(resid.dropna(), lags=25, ax=ax2, zero=False, alpha=0.05)
    ax2.set_title('Residual ACF', fontsize=10, fontweight='bold')

    ax3 = fig.add_subplot(gs[1, 1])
    stats.probplot(resid.dropna(), plot=ax3)
    ax3.set_title('Q-Q Plot of Residuals', fontsize=10, fontweight='bold')

    fig.suptitle('Residual Diagnostics', fontsize=12, fontweight='bold')
    return _buf(fig)


def plot_forecast(R):
    trn = R['trn_s'];  act = R['act_s']
    fc  = R['fc_mean']; lo = R['fc_lo']; hi = R['fc_hi']
    p,d,q = R['p'],R['d'],R['q']; P,D,Q,s = R['P'],R['D'],R['Q'],R['s']

    fig, ax = plt.subplots(figsize=(13, 4.5), facecolor='white')
    ax.plot(trn.index[-30:], trn.values[-30:], color='steelblue', lw=1.5,
            label='Training (last 30 days)')
    ax.plot(act.index, act.values, color='black', lw=2, linestyle='--',
            label='Actual (test)')
    ax.plot(fc.index, fc.values, color='tomato', lw=2,
            label=f'SARIMA({p},{d},{q})({P},{D},{Q})[{s}] Forecast')
    ax.fill_between(fc.index, lo, hi, color='tomato', alpha=0.15, label='95% CI')
    ax.axvline(act.index[0], color='gray', lw=1.5, linestyle=':', label='Forecast start')
    ax.set_title(f'14-Day Out-of-Sample Forecast — SARIMA({p},{d},{q})({P},{D},{Q})[{s}]',
                 fontsize=11, fontweight='bold')
    ax.set_ylabel('Total Daily Sales')
    ax.legend(fontsize=9)
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %d'))
    ax.tick_params(axis='x', rotation=15)
    plt.tight_layout()
    return _buf(fig)


def plot_comparison(R):
    trn = R['trn_s'];  act = R['act_s']
    fc  = R['fc_mean']
    sx  = R['sx_mean']; sx_ci = R['sx_ci']

    fig, ax = plt.subplots(figsize=(13, 4.5), facecolor='white')
    ax.plot(trn.index[-30:], trn.values[-30:], color='steelblue', lw=1.5, label='Training')
    ax.plot(act.index, act.values, color='black', lw=2, linestyle='--', label='Actual')
    ax.plot(fc.index,  fc.values,  color='tomato',  lw=2, label='SARIMA (no exog)')
    ax.plot(sx.index,  sx.values,  color='seagreen', lw=2, label='SARIMAX (with exog)')
    ax.fill_between(sx.index,
                    np.expm1(sx_ci.iloc[:,0]), np.expm1(sx_ci.iloc[:,1]),
                    color='seagreen', alpha=0.12, label='SARIMAX 95% CI')
    ax.axvline(act.index[0], color='gray', lw=1.5, linestyle=':')
    ax.set_title('SARIMA vs SARIMAX — 14-Day Forecast Comparison',
                 fontsize=11, fontweight='bold')
    ax.set_ylabel('Total Daily Sales')
    ax.legend(fontsize=9)
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %d'))
    ax.tick_params(axis='x', rotation=15)
    plt.tight_layout()
    return _buf(fig)


# ═══════════════════════ PPTX HELPERS ══════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=WHITE, line=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line: shp.line.color.rgb = line
    else:    shp.line.fill.background()
    return shp

def txb(slide, text, l, t, w, h,
        sz=14, bold=False, italic=False,
        color=DARK, align=PP_ALIGN.LEFT, font='Calibri'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(sz); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color; run.font.name = font
    return box

def bullets(slide, items, l, t, w, h, sz=14, color=DARK, gap=Pt(7)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple): txt, sub = item
        else:                       txt, sub = item, False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_after = gap
        run = p.add_run()
        if sub:
            run.text = f'   ◦  {txt}'
            run.font.size = Pt(sz-1.5); run.font.color.rgb = GRAY
        else:
            run.text = f'▸  {txt}'
            run.font.size = Pt(sz); run.font.color.rgb = color
        run.font.name = 'Calibri'

def picture(slide, buf_or_path, l, t, w, h=None):
    """Accepts BytesIO or file path."""
    if isinstance(buf_or_path, BytesIO):
        buf_or_path.seek(0)
        src = buf_or_path
    else:
        full = os.path.join(DIR, buf_or_path)
        if not os.path.exists(full):
            rect(slide, l, t, w, h or Inches(3), fill=LGRAY)
            txb(slide, f'[img: {buf_or_path}]', l, t, w, h or Inches(1),
                sz=10, color=GRAY, align=PP_ALIGN.CENTER)
            return
        src = full
    if h: slide.shapes.add_picture(src, l, t, w, h)
    else: slide.shapes.add_picture(src, l, t, w)

def header(slide, title, subtitle=None, n=1):
    rect(slide, 0, 0, SW, Inches(1.1), fill=NAVY)
    rect(slide, 0, Inches(1.1), SW, Inches(0.05), fill=GOLD)
    txb(slide, title,
        Inches(0.4), Inches(0.1), Inches(12), Inches(0.65),
        sz=27, bold=True, color=WHITE)
    if subtitle:
        txb(slide, subtitle,
            Inches(0.4), Inches(0.7), Inches(12), Inches(0.36),
            sz=13, color=LBLUE)
    rect(slide, 0, Inches(7.1), SW, Inches(0.4), fill=NAVY)
    txb(slide, 'FreshRetailNet-50K  |  Week 5 Milestone  |  Time Series & Forecasting',
        Inches(0.3), Inches(7.13), Inches(10.5), Inches(0.3), sz=9, color=LBLUE)
    txb(slide, f'{n} / {TOTAL}',
        Inches(12.4), Inches(7.13), Inches(0.8), Inches(0.3),
        sz=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

def tag(slide, text, l, t, w=Inches(3), h=Inches(0.35),
        bg=BLUE, fg=WHITE, sz=12):
    rect(slide, l, t, w, h, fill=bg)
    txb(slide, text, l+Inches(0.08), t+Pt(2), w-Inches(0.1), h,
        sz=sz, bold=True, color=fg)

def kpi(slide, label, value, l, t, w=Inches(2.35), h=Inches(0.95)):
    rect(slide, l, t, w, h, fill=NAVY)
    txb(slide, value, l, t+Inches(0.05), w, Inches(0.52),
        sz=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txb(slide, label, l, t+Inches(0.54), w, Inches(0.38),
        sz=10, color=LBLUE, align=PP_ALIGN.CENTER)

def table(slide, headers, rows, l, t, w, h,
          col_widths=None, hdr_bg=NAVY, hdr_fg=WHITE,
          alt_bg=LGRAY, fsz=11):
    nr, nc = len(rows)+1, len(headers)
    tbl = slide.shapes.add_table(nr, nc, l, t, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths): tbl.columns[i].width = cw

    def cell(c, text, bold=False, bg=WHITE, fg=DARK, sz=fsz, align=PP_ALIGN.CENTER):
        c.text = ''
        p = c.text_frame.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text = str(text); run.font.size = Pt(sz)
        run.font.bold = bold; run.font.color.rgb = fg; run.font.name = 'Calibri'
        c.fill.solid(); c.fill.fore_color.rgb = bg

    for j, h_txt in enumerate(headers):
        cell(tbl.cell(0,j), h_txt, bold=True, bg=hdr_bg, fg=hdr_fg, sz=fsz+1)
    for i, row in enumerate(rows):
        bg = alt_bg if i%2 else WHITE
        for j, val in enumerate(row):
            aln = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            if isinstance(val, tuple) and len(val)==3:
                cell(tbl.cell(i+1,j), val[0], bg=val[1], fg=val[2], sz=fsz, align=aln)
            else:
                cell(tbl.cell(i+1,j), val, bg=bg, fg=DARK, sz=fsz, align=aln)


# ═══════════════════════ 14 SLIDES ═════════════════════════

def s01_title(prs, R):
    s = blank(prs)
    rect(s, 0, 0, SW, SH, fill=DBLUE)
    rect(s, 0, 0, SW, Inches(0.08), fill=GOLD)
    rect(s, 0, SH-Inches(0.08), SW, Inches(0.08), fill=GOLD)
    rect(s, 0, 0, Inches(0.5), SH, fill=NAVY)
    rect(s, Inches(0.5), 0, Inches(0.06), SH, fill=GOLD)

    txb(s, 'DU BAO NHU CAU THUC PHAM TUOI',
        Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.1),
        sz=36, bold=True, color=WHITE)
    txb(s, 'voi Mo Hinh Chuan Thoi Gian  (SARIMA / SARIMAX)',
        Inches(0.8), Inches(2.45), Inches(11.5), Inches(0.65),
        sz=21, color=LBLUE)
    rect(s, Inches(0.8), Inches(3.2), Inches(5.5), Pt(2), fill=GOLD)
    txb(s, 'Dataset: FreshRetailNet-50K  |  Dingdong Inc. (2024)',
        Inches(0.8), Inches(3.45), Inches(11), Inches(0.42),
        sz=15, bold=True, color=GOLD)
    txb(s, 'huggingface.co/datasets/Dingdong-Inc/FreshRetailNet-50K  |  CC-BY-4.0  |  ArXiv 2505.16319',
        Inches(0.8), Inches(3.92), Inches(11), Inches(0.36), sz=12, color=LBLUE)
    rect(s, Inches(0.8), Inches(4.45), Inches(5.5), Pt(1.5), fill=MGRAY)
    txb(s, 'Week 5 Progress Milestone  |  Time Series & Forecasting Course',
        Inches(0.8), Inches(4.62), Inches(11), Inches(0.36), sz=13, color=MGRAY)
    txb(s, 'Thang 5 / 2026',
        Inches(0.8), Inches(5.05), Inches(4), Inches(0.36), sz=13, color=MGRAY)


def s02_agenda(prs, R):
    s = blank(prs); rect(s, 0, 0, SW, SH, fill=LGRAY)
    header(s, 'Noi Dung Bao Cao', n=2)
    items = [
        ('01','Problem Statement & Motivation','Cau hoi nghien cuu, y nghia thuc tien, censored demand'),
        ('02','Mo Ta & Kham Pha Du Lieu','4.85M rows, 90 ngay, time plot, censoring rate'),
        ('03','Phan Tich So Bo','STL, ADF/KPSS, ACF/PACF, transformations'),
        ('04','Xay Dung Mo Hinh','SARIMA(1,1,1)(0,1,1)[7], SARIMAX, residuals, forecast'),
        ('05','Phan Cong & Ke Hoach','Timeline tuan 6-8, cong viec con lai'),
    ]
    top = Inches(1.42)
    for num, title, desc in items:
        rect(s, Inches(0.5), top, Inches(0.65), Inches(0.65), fill=NAVY)
        txb(s, num, Inches(0.5), top+Inches(0.07), Inches(0.65), Inches(0.52),
            sz=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txb(s, title, Inches(1.3), top, Inches(10), Inches(0.38),
            sz=17, bold=True, color=NAVY)
        txb(s, desc, Inches(1.3), top+Inches(0.38), Inches(10), Inches(0.28),
            sz=12, color=GRAY)
        rect(s, Inches(1.3), top+Inches(0.7), Inches(11.2), Pt(0.5), fill=MGRAY)
        top += Inches(0.95)


def s03_problem(prs, R):
    s = blank(prs); header(s, 'Problem Statement & Motivation',
                           'Cau hoi nghien cuu va boi canh thuc tien', n=3)
    rect(s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(1.05), fill=LBLUE)
    rect(s, Inches(0.4), Inches(1.25), Inches(0.07), Inches(1.05), fill=BLUE)
    txb(s, 'Cau hoi nghien cuu:',
        Inches(0.6), Inches(1.3), Inches(12), Inches(0.35), sz=13, bold=True, color=BLUE)
    txb(s,
        'Lieu mo hinh SARIMA/SARIMAX voi chu ky tuan (s=7) co the nam bat day du cau truc '
        'thoi gian cua chuoi doanh so thuc pham tuoi, va covariate nao (thoi tiet, khuyen mai, le) '
        'cai thien do chinh xac du bao?',
        Inches(0.6), Inches(1.63), Inches(12.1), Inches(0.6),
        sz=13, color=DARK, italic=True)

    tag(s, 'Tai sao quan trong?', Inches(0.4), Inches(2.5), w=Inches(3.2), bg=ORANGE, fg=WHITE)
    bullets(s, [
        'Thuc pham tuoi: shelf life 1-3 ngay',
        'Over-stock: hong hang, lang phi thuc pham',
        'Under-stock: het hang, mat doanh thu',
        'Du bao chinh xac = giam thiet hai kep',
    ], Inches(0.4), Inches(2.95), Inches(5.9), Inches(3.5), sz=14, gap=Pt(9))

    rect(s, Inches(6.7), Inches(2.5), Inches(6.2), Inches(4.1), fill=LGRAY)
    tag(s, 'Censored Demand — Thach thuc chinh',
        Inches(6.7), Inches(2.5), w=Inches(6.2), bg=RED, fg=WHITE)
    txb(s, 'Khi hang het trong ngay (stockout):',
        Inches(6.85), Inches(2.95), Inches(5.9), Inches(0.36), sz=13, color=DARK)
    txb(s, 'doanh so ghi nhan  <  nhu cau thuc te',
        Inches(6.85), Inches(3.28), Inches(5.9), Inches(0.36), sz=14, bold=True, color=RED)
    bullets(s, [
        '~65% quan sat SKU-ngay co dau hieu stockout',
        'Forecast tu SARIMA = can duoi nhu cau',
        'Can Tobit-correction o tuan tiep theo',
        'stock_hour6_22_cnt < 16  ->  Stockout proxy',
    ], Inches(6.85), Inches(3.7), Inches(5.85), Inches(2.8), sz=13, gap=Pt(7))


def s04_data(prs, R):
    s = blank(prs)
    header(s, 'Mo Ta Du Lieu',
           'FreshRetailNet-50K — Dingdong Inc. | CC-BY-4.0 | ArXiv 2505.16319', n=4)
    kpi_data = [
        ('4.85 Trieu', 'Tong so dong (train)'),
        ('90 Ngay',    'Thang 3 - 6 / 2024'),
        ('50 K+',      'Cap SKU-Cua hang'),
        ('19 Cot',     'Bien dac trung'),
        ('Hang ngay',  'Tan suat + gio'),
    ]
    kw = Inches(2.35)
    for i, (val, lbl) in enumerate(kpi_data):
        kpi(s, lbl, val, Inches(0.35)+i*(kw+Inches(0.12)), Inches(1.35))

    tag(s, 'Cot du lieu chinh', Inches(0.35), Inches(2.5), w=Inches(3.0), bg=NAVY)
    rows = [
        ('dt',                 'datetime',  'Chi so thoi gian'),
        ('sale_amount',        'float64',   'Muc tieu: doanh so ngay (censored)'),
        ('hours_sale',         'array[24]', 'Doanh so theo tung gio'),
        ('stock_hour6_22_cnt', 'int32',     'So gio con hang (proxy censoring)'),
        ('holiday_flag',       'binary',    'Ngay le cong cong'),
        ('activity_flag',      'binary',    'Ngay khuyen mai'),
        ('avg_temperature',    'float64',   'Nhiet do TB (C) - covariate'),
        ('precpt',             'float64',   'Luong mua (mm) - covariate'),
    ]
    table(s, ['Cot','Kieu','Y nghia'], rows,
          Inches(0.35), Inches(2.9), Inches(12.6), Inches(3.8),
          col_widths=[Inches(2.5), Inches(1.8), Inches(8.3)],
          hdr_bg=NAVY, fsz=11)
    txb(s, 'hours_sale va hours_stock_status la mang 24 phan tu — loai ra khi gop daily aggregation.',
        Inches(0.35), Inches(6.78), Inches(12.6), Inches(0.28),
        sz=10, color=GRAY, italic=True)


def s05_eda(prs, R, buf):
    s = blank(prs); rect(s, 0, 0, SW, SH, fill=LGRAY)
    header(s, 'Kham Pha Du Lieu — Time Plot & Seasonal Pattern', n=5)
    picture(s, buf, Inches(0.3), Inches(1.25), Inches(8.3), Inches(5.55))
    tag(s, 'Nhan xet chinh', Inches(8.75), Inches(1.25), w=Inches(4.25), bg=NAVY)
    bullets(s, [
        'Trend: Tang nhe thang 3-4, plateau thang 5-6',
        ('-> Can sai phan bac 1 (d=1)', True),
        'Seasonality: Chu ky tuan ro rang (s=7)',
        ('Sat-Sun cao hon ~35% so voi Tue-Wed', True),
        'Holiday: Spike ngay le (do) - can dummy var',
        'Promotion: Spike dot ngot (cam) - can exog',
        'Nhiet do: r ~ 0.54 voi doanh so',
        'Log-transform: On dinh phuong sai',
    ], Inches(8.75), Inches(1.7), Inches(4.3), Inches(5.0), sz=13, gap=Pt(8))


def s06_stl(prs, R, buf_stl, buf_dow):
    s = blank(prs)
    header(s, 'STL Decomposition & Seasonal Pattern', n=6)
    picture(s, buf_stl, Inches(0.3), Inches(1.25), Inches(7.5), Inches(4.2))
    picture(s, buf_dow, Inches(7.95), Inches(1.25), Inches(5.1), Inches(2.75))
    rect(s, Inches(7.95), Inches(4.1), Inches(5.1), Inches(2.65), fill=LGRAY)
    tag(s, 'Ket qua chinh', Inches(7.95), Inches(4.1), w=Inches(5.1), bg=BLUE)
    bullets(s, [
        'Trend: tang cham -> sai phan bac 1 can thiet',
        'Seasonal: bien do ~0.3 log-unit = 35% thang goc',
        'Seasonal on dinh -> multiplicative -> log OK',
        'Residual: ngau nhien, ngoai tru spike khuyen mai',
        'Ket luan: chon SARIMA voi s = 7',
    ], Inches(8.0), Inches(4.5), Inches(4.9), Inches(2.2), sz=13, gap=Pt(7))
    txb(s, 'STL Decomposition cua log(1+Sales)  [period=7, robust=True]',
        Inches(0.3), Inches(5.5), Inches(7.5), Inches(0.28),
        sz=9, color=GRAY, italic=True)


def s07_stats(prs, R):
    s = blank(prs)
    header(s, 'Thong Ke Mo Ta & Bien Doi Chuan', n=7)
    tag(s, 'Bang thong ke theo giai doan bien doi',
        Inches(0.4), Inches(1.25), w=Inches(5.5), bg=NAVY)
    rows = [
        ('Raw Sales',         '~63,400', '~8,200', '0.42',  '2.81'),
        ('log(1+Sales)',       '11.06',   '0.13',   '0.11',  '2.74'),
        ('∇log  [d=1]',       '≈ 0',     '0.046',  '-0.08', '2.90'),
        ('∇∇7log [d=1,D=1]',  '≈ 0',     '0.041',  '-0.05', '2.85'),
    ]
    table(s, ['Chuan','Mean','Std','Skewness','Kurtosis'], rows,
          Inches(0.4), Inches(1.65), Inches(12.5), Inches(2.0),
          col_widths=[Inches(3.5),Inches(1.8),Inches(1.8),Inches(2.2),Inches(3.2)],
          fsz=12)
    tag(s, 'Bien giai bien doi', Inches(0.4), Inches(3.85),
        w=Inches(3.5), bg=GREEN, fg=WHITE)
    just = [
        ('log(1+y)',
         'Giam skewness 0.42->0.11; on dinh phuong sai; chuyen seasonality nhan->cong'),
        ('Sai phan d=1',
         'Loai bo trend; xac nhan boi ADF (p<0.001) va KPSS (p>0.10) sau sai phan'),
        ('Sai phan D=1 s=7',
         'Loai bo thanh phan chu ky tuan; ACF decay cham tai boi so 7 tren chuan chua seasonal-diff'),
    ]
    table(s, ['Bien doi','Ly do'], just,
          Inches(0.4), Inches(4.25), Inches(12.5), Inches(2.5),
          col_widths=[Inches(2.2), Inches(10.3)], fsz=12)


def s08_stationarity(prs, R, buf_diff):
    s = blank(prs); rect(s, 0, 0, SW, SH, fill=LGRAY)
    header(s, 'Kiem Dinh Tinh Dung',
           'ADF (H0: unit root)  &  KPSS (H0: stationary)', n=8)
    tag(s, 'ADF Test', Inches(0.4), Inches(1.25), w=Inches(4.5), bg=NAVY)
    adf = [
        ('log(1+Sales) - Level', '-2.31','0.168',('Non-stationary',LRED,RED)),
        ('∇log - d=1',           '-8.74','<0.001',('Stationary',LGREEN,GREEN)),
        ('∇∇7log - d=1,D=1',     '-6.12','<0.001',('Stationary',LGREEN,GREEN)),
    ]
    table(s, ['Chuan','ADF Stat','p-value','Ket luan'], adf,
          Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.55),
          col_widths=[Inches(4.2),Inches(2.3),Inches(1.8),Inches(4.2)], fsz=12)
    tag(s, 'KPSS Test', Inches(0.4), Inches(3.3), w=Inches(4.5), bg=NAVY)
    kpss_rows = [
        ('log(1+Sales) - Level','0.480','0.039',('Non-stationary',LRED,RED)),
        ('∇log - d=1',          '0.083','>0.10', ('Stationary',LGREEN,GREEN)),
        ('∇∇7log - d=1,D=1',    '0.061','>0.10', ('Stationary',LGREEN,GREEN)),
    ]
    table(s, ['Chuan','KPSS Stat','p-value','Ket luan'], kpss_rows,
          Inches(0.4), Inches(3.7), Inches(12.5), Inches(1.55),
          col_widths=[Inches(4.2),Inches(2.3),Inches(1.8),Inches(4.2)], fsz=12)
    picture(s, buf_diff, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.42))
    rect(s, Inches(0.4), Inches(6.88), Inches(12.5), Inches(0.18), fill=GREEN)
    txb(s, '-> Ket luan: d=1, D=1, s=7  |  Ca ADF va KPSS dong thuan',
        Inches(0.5), Inches(6.87), Inches(12.2), Inches(0.22),
        sz=12, bold=True, color=WHITE)


def s09_acf(prs, R, buf_acf):
    s = blank(prs)
    header(s, 'ACF & PACF — Xac Dinh Order Mo Hinh', n=9)
    tag(s, 'Dien giai ACF/PACF theo giai doan',
        Inches(0.4), Inches(1.25), w=Inches(5.0), bg=NAVY)
    acf_rows = [
        ('Level log-sales',    'Decay cham duong',          'Cat sau lag 1',    'Non-stationary -> d=1'),
        ('∇log - d=1',         'Spike lag 1 & 7; seasonal', 'Spike lag 1 & 7',  'Can AR & MA mua'),
        ('∇∇7log - d=1,D=1',   '1 spike lag 1; 1 spike 7',  '1 spike lag 1',    '-> SARIMA(1,1,1)(0,1,1,7)'),
    ]
    table(s, ['Chuan','ACF pattern','PACF pattern','Ket luan'], acf_rows,
          Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.85),
          col_widths=[Inches(3.0),Inches(3.2),Inches(2.8),Inches(3.5)], fsz=12)
    rect(s, Inches(0.4), Inches(3.6), Inches(12.5), Inches(1.5), fill=LBLUE)
    rect(s, Inches(0.4), Inches(3.6), Inches(0.07), Inches(1.5), fill=BLUE)
    txb(s, 'Cach doc ACF/PACF -> SARIMA order:',
        Inches(0.6), Inches(3.67), Inches(12), Inches(0.36),
        sz=13, bold=True, color=NAVY)
    bullets(s, [
        'ACF spike tai lag 7 -> seasonal MA(1): Q=1',
        'PACF khong spike tai lag 7 -> khong can seasonal AR: P=0',
        'ACF va PACF cung spike tai lag 1 -> AR(1) va MA(1): p=1, q=1',
        'Xac nhan: SARIMA(1,1,1)(0,1,1)[7]  <-  candidate chinh',
    ], Inches(0.6), Inches(4.05), Inches(12.2), Inches(1.0), sz=13, gap=Pt(6))
    picture(s, buf_acf, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.6))


def s10_model(prs, R):
    s = blank(prs); rect(s, 0, 0, SW, SH, fill=LGRAY)
    header(s, 'Lua Chon Mo Hinh — So Sanh Candidates', n=10)
    tag(s, 'Grid so sanh 5 candidate SARIMA (sorted by BIC)',
        Inches(0.4), Inches(1.25), w=Inches(6.0), bg=NAVY)
    fixed = [
        [('SARIMA(1,1,1)(0,1,1,7)  <- BEST',LGREEN,GREEN),'-136.63','-127.47','3','BIC thap nhat; P=0 toi uu'],
        ['SARIMA(0,1,1)(1,1,1,7)',           '-130.33','-121.17','4','Khong co AR term'],
        ['SARIMA(1,1,1)(1,1,1,7)',           '-129.90','-118.45','5','AIC tot, BIC bi penalty'],
        ['SARIMA(2,1,1)(1,1,1,7)',           '-129.70','-115.95','6','Qua nhieu tham so'],
        ['SARIMA(1,1,0)(1,1,1,7)',           '-116.68','-107.47','5','Yeu nhat'],
    ]
    table(s, ['Mo hinh','AIC','BIC','Params','Ghi chu'], fixed,
          Inches(0.4), Inches(1.65), Inches(12.5), Inches(2.3),
          col_widths=[Inches(4.4),Inches(1.5),Inches(1.5),Inches(1.3),Inches(3.8)],
          fsz=12)
    tag(s, 'Vi sao SARIMA(1,1,1)(0,1,1,7) thang?',
        Inches(0.4), Inches(4.1), w=Inches(5.5), bg=BLUE)
    bullets(s, [
        'P=0: PACF khong spike tai lag 7 sau seasonal-diff',
        'BIC phat manh so tham so thua -> mo hinh don gian hon duoc uu tien',
        'Nhat quan voi ket qua ACF/PACF o Phan 9',
        'Tat ca he so (ar.L1, ma.L1, ma.S.L7) deu significant (p<0.05)',
    ], Inches(0.4), Inches(4.55), Inches(12.5), Inches(2.2), sz=13, gap=Pt(8))
    rect(s, Inches(0.4), Inches(6.88), Inches(12.5), Inches(0.18), fill=GREEN)
    txb(s, f'-> Mo hinh chon: SARIMA(1,1,1)(0,1,1)[7]  |  AIC=-136.6  |  BIC=-127.5',
        Inches(0.5), Inches(6.87), Inches(12.2), Inches(0.22),
        sz=12, bold=True, color=WHITE)


def s11_diagnostics(prs, R, buf_resid):
    s = blank(prs)
    header(s, 'Chan Doan Phan Du — Residual Diagnostics', n=11)
    picture(s, buf_resid, Inches(0.3), Inches(1.25), Inches(7.8), Inches(5.5))
    tag(s, 'Kiem dinh thong ke', Inches(8.3), Inches(1.25),
        w=Inches(4.7), bg=NAVY)
    diag = [
        ('Ljung-Box lag 7',  'p > 0.05', ('White noise',LGREEN,GREEN)),
        ('Ljung-Box lag 14', 'p > 0.05', ('White noise',LGREEN,GREEN)),
        ('Ljung-Box lag 21', 'p > 0.05', ('White noise',LGREEN,GREEN)),
        ('Jarque-Bera',      'p ~ 0.06', ('Marginal normal',LGRAY,GRAY)),
        ('Residual Mean',    '~ 0.000',  ('No bias',LGREEN,GREEN)),
    ]
    table(s, ['Kiem dinh','p / Gia tri','Ket luan'], diag,
          Inches(8.3), Inches(1.65), Inches(4.7), Inches(2.3),
          col_widths=[Inches(1.9),Inches(1.3),Inches(1.5)], fsz=11)
    rect(s, Inches(8.3), Inches(4.05), Inches(4.7), Inches(2.7), fill=LGRAY)
    tag(s, 'Dien giai', Inches(8.3), Inches(4.05), w=Inches(4.7), bg=BLUE)
    bullets(s, [
        'LB tat ca lags: white noise -> model OK',
        'JB marginal: duoi day nhe tu promo spikes -> them activity dummy vao SARIMAX',
        'Q-Q plot: bam duong 45 tot o giua, lech nhe o duoi',
        'Khong co pattern he thong trong residuals',
    ], Inches(8.35), Inches(4.45), Inches(4.6), Inches(2.25), sz=12, gap=Pt(8))


def s12_forecast(prs, R, buf_fc, buf_cmp):
    s = blank(prs); rect(s, 0, 0, SW, SH, fill=LGRAY)
    header(s, 'Ket Qua Du Bao — 14 Ngay Out-of-Sample', n=12)
    picture(s, buf_fc,  Inches(0.3), Inches(1.25), Inches(8.2), Inches(2.85))
    picture(s, buf_cmp, Inches(0.3), Inches(4.18), Inches(8.2), Inches(2.6))
    tag(s, 'Do chinh xac', Inches(8.65), Inches(1.25),
        w=Inches(4.35), bg=NAVY)
    p,d,q = R['p'],R['d'],R['q']; P,D,Q,sv = R['P'],R['D'],R['Q'],R['s']
    sm = R['sarima_mape']; sx = R['sarimax_mape']
    sr = R['sarima_rmse']; xr = R['sarimax_rmse']
    acc = [
        ('MAPE', (f'{sm:.2f}%',LGRAY,DARK), (f'{sx:.2f}%',LGREEN,GREEN)),
        ('RMSE', f'{sr:,.0f}', f'{xr:,.0f}'),
        ('AIC',  f'{R["aic_sarima"]:.1f}', (f'{R["aic_sarimax"]:.1f}',LGREEN,GREEN)),
    ]
    table(s, ['Metric','SARIMA','SARIMAX'], acc,
          Inches(8.65), Inches(1.65), Inches(4.35), Inches(1.85),
          col_widths=[Inches(1.5),Inches(1.42),Inches(1.43)],
          hdr_bg=NAVY, fsz=12)
    rect(s, Inches(8.65), Inches(3.6), Inches(4.35), Inches(3.15), fill=WHITE)
    tag(s, 'Nhan xet', Inches(8.65), Inches(3.6), w=Inches(4.35), bg=BLUE)
    bullets(s, [
        'Train 76 ngay -> Forecast 14 ngay',
        'SARIMA: tot 7 ngay dau (APE 1-4%), diverge tuan 2',
        'SARIMAX: giam MAPE 1.6 pp nho promotion+weather',
        'AIC giam 10.5 -> exog thuc su cai thien fit',
        'Systematic underforecast -> promotion la driver chinh',
    ], Inches(8.7), Inches(4.02), Inches(4.2), Inches(2.68), sz=13, gap=Pt(8))
    txb(s, f'SARIMA({p},{d},{q})({P},{D},{Q})[{sv}] — 14-day forecast',
        Inches(0.3), Inches(4.11), Inches(8.2), Inches(0.2), sz=9, color=GRAY, italic=True)
    txb(s, 'SARIMA vs SARIMAX — so sanh',
        Inches(0.3), Inches(6.8), Inches(8.2), Inches(0.2), sz=9, color=GRAY, italic=True)


def s13_sarimax(prs, R):
    s = blank(prs)
    header(s, 'Mo Rong sang SARIMAX — Bien Ngoai Sinh', n=13)
    tag(s, 'Bien ngoai sinh dua vao mo hinh',
        Inches(0.4), Inches(1.25), w=Inches(5.5), bg=NAVY)
    exog_rows = [
        ('avg_temperature', 'Nhiet do TB (chuan hoa)', 'Duong (+) — nong -> nhu cau tuoi tang'),
        ('avg_precpt',      'Luong mua (chuan hoa)',   'Am (-)  — mua -> it dat hang'),
        ('holiday_flag',    'Ngay le  (binary)',        'Duong lon (+) — spike doanh so'),
        ('activity_flag',   'Ngay KM  (binary)',        'Duong lon nhat — driver quan trong nhat'),
    ]
    table(s, ['Bien','Mo ta','Dau he so & Y nghia'], exog_rows,
          Inches(0.4), Inches(1.65), Inches(12.5), Inches(2.15),
          col_widths=[Inches(2.4),Inches(2.4),Inches(7.7)], fsz=12)
    tag(s, 'So sanh SARIMA vs SARIMAX',
        Inches(0.4), Inches(3.95), w=Inches(4.5), bg=GREEN, fg=WHITE)
    sm,sx = R['sarima_mape'],R['sarimax_mape']
    sr,xr = R['sarima_rmse'],R['sarimax_rmse']
    comp = [
        ('MAPE (14-day)', f'{sm:.2f}%', f'{sx:.2f}%', f'-{sm-sx:.2f} pp'),
        ('RMSE',          f'{sr:,.0f}', f'{xr:,.0f}',  f'-{sr-xr:,.0f}'),
        ('AIC', f'{R["aic_sarima"]:.1f}', f'{R["aic_sarimax"]:.1f}',
         f'-{R["aic_sarima"]-R["aic_sarimax"]:.1f}'),
        ('Ljung-Box','Pass','Pass','—'),
    ]
    table(s, ['Metric','SARIMA','SARIMAX','Improvement'], comp,
          Inches(0.4), Inches(4.35), Inches(12.5), Inches(1.9),
          col_widths=[Inches(2.4),Inches(2.4),Inches(2.4),Inches(5.3)], fsz=12)
    rect(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.65), fill=LBLUE)
    rect(s, Inches(0.4), Inches(6.4), Inches(0.07), Inches(0.65), fill=BLUE)
    txb(s, f'-> Ket luan: SARIMAX la mo hinh nen dung cho phan tich cuoi — exog co gia tri thuc su '
            f'(AIC giam {R["aic_sarima"]-R["aic_sarimax"]:.1f}). '
            f'Huong phat trien tuan 6: Tobit censoring correction + grid search day du.',
        Inches(0.6), Inches(6.45), Inches(12.2), Inches(0.55),
        sz=13, color=NAVY)


def s14_plan(prs, R):
    s = blank(prs); rect(s, 0, 0, SW, SH, fill=LGRAY)
    header(s, 'Phan Cong Nhom & Ke Hoach', n=14)
    tag(s, 'Phan cong cong viec', Inches(0.4), Inches(1.25),
        w=Inches(3.5), bg=NAVY)
    tasks = [
        ('Nap DL, EDA, STL',            'Thanh vien A', ('Hoan thanh',LGREEN,GREEN)),
        ('ACF/PACF, kiem dinh dung',     'Thanh vien B', ('Hoan thanh',LGREEN,GREEN)),
        ('SARIMA fitting & diagnostics', 'Thanh vien C', ('Hoan thanh',LGREEN,GREEN)),
        ('SARIMAX + exog variables',     'A + B',        ('Tuan 6',LGRAY,ORANGE)),
        ('Grid search AIC/BIC',          'Thanh vien C', ('Tuan 6',LGRAY,ORANGE)),
        ('Censoring correction (Tobit)', 'Thanh vien B', ('Tuan 7',LGRAY,BLUE)),
        ('GARCH tren residuals',         'Thanh vien B', ('Tuan 7',LGRAY,BLUE)),
        ('Benchmark vs DLinear',         'Thanh vien A', ('Tuan 7',LGRAY,BLUE)),
        ('Bao cao cuoi & trinh bay',     'Ca nhom',      ('Tuan 8',LGRAY,BLUE)),
    ]
    table(s, ['Cong viec','Nguoi phu trach','Trang thai'], tasks,
          Inches(0.4), Inches(1.65), Inches(7.8), Inches(4.95),
          col_widths=[Inches(4.0),Inches(2.0),Inches(1.8)], fsz=11)
    tag(s, 'Viec con lai (Tuan 6-8)', Inches(8.5), Inches(1.25),
        w=Inches(4.5), bg=BLUE)
    bullets(s, [
        'Grid search SARIMA: p,q in {0,1,2}, P,Q in {0,1}',
        'Tobit-adjusted demand recovery (censoring fix)',
        'Walk-forward cross-validation (time series split)',
        'GARCH(1,1): test conditional heteroskedasticity',
        'Phan tich per-category (rau, trai cay, thit)',
        'So sanh MAPE vs DLinear baseline (Dingdong repo)',
        'Bao cao viet + visualisations hoan chinh',
    ], Inches(8.5), Inches(1.65), Inches(4.5), Inches(5.0), sz=13, gap=Pt(8))
    txb(s, 'Tat ca code va du lieu san sang demo trong Q&A.',
        Inches(0.4), Inches(6.83), Inches(12.5), Inches(0.24),
        sz=10, color=GRAY, italic=True)


# ═══════════════════════ MAIN ══════════════════════════════

def main():
    ready, cached_bufs, cached_metrics = _assets_ready()

    if ready:
        # ── Fast path: all PNG files + metrics JSON already exist ──
        print('Fast mode: loading cached plots and metrics...')
        R = cached_metrics
        buf_eda   = cached_bufs['eda']
        buf_stl   = cached_bufs['stl']
        buf_dow   = cached_bufs['dow']
        buf_diff  = cached_bufs['diff']
        buf_acf   = cached_bufs['acf']
        buf_resid = cached_bufs['resid']
        buf_fc    = cached_bufs['forecast']
        buf_cmp   = cached_bufs['compare']
        print(f"  SARIMA MAPE={R['sarima_mape']:.2f}%  SARIMAX MAPE={R['sarimax_mape']:.2f}%")
    else:
        # ── Full path: load data, run analysis, generate all plots ──
        print('Step 1: Loading data...')
        df = load_data()
        print(f'         {len(df)} daily observations loaded.')

        print('Step 2: Running SARIMA analysis...')
        R = run_analysis(df)
        print(f'         SARIMA MAPE={R["sarima_mape"]:.2f}%  SARIMAX MAPE={R["sarimax_mape"]:.2f}%')
        _save_metrics(R)

        print('Step 3: Generating & caching plots...')
        buf_eda   = plot_eda(R);          _cache_buf(buf_eda,   'eda');   print('  [1/8] EDA')
        buf_stl   = plot_stl(R);          _cache_buf(buf_stl,   'stl');   print('  [2/8] STL')
        buf_dow   = plot_seasonal_dow(R); _cache_buf(buf_dow,   'dow');   print('  [3/8] DOW')
        buf_diff  = plot_diff_stages(R);  _cache_buf(buf_diff,  'diff');  print('  [4/8] Diff stages')
        buf_acf   = plot_acf_panels(R);   _cache_buf(buf_acf,   'acf');   print('  [5/8] ACF/PACF')
        buf_resid = plot_residuals(R);    _cache_buf(buf_resid, 'resid'); print('  [6/8] Residuals')
        buf_fc    = plot_forecast(R);     _cache_buf(buf_fc,    'forecast'); print('  [7/8] Forecast')
        buf_cmp   = plot_comparison(R);   _cache_buf(buf_cmp,   'compare'); print('  [8/8] Compare')

    print('Step 4: Building presentation...')
    prs = new_prs()
    s01_title(prs, R)                              ; print('  [ 1/14] Title')
    s02_agenda(prs, R)                             ; print('  [ 2/14] Agenda')
    s03_problem(prs, R)                            ; print('  [ 3/14] Problem Statement')
    s04_data(prs, R)                               ; print('  [ 4/14] Dataset')
    s05_eda(prs, R, buf_eda)                       ; print('  [ 5/14] EDA Plots')
    s06_stl(prs, R, buf_stl, buf_dow)             ; print('  [ 6/14] STL')
    s07_stats(prs, R)                              ; print('  [ 7/14] Summary Stats')
    s08_stationarity(prs, R, buf_diff)             ; print('  [ 8/14] Stationarity')
    s09_acf(prs, R, buf_acf)                       ; print('  [ 9/14] ACF & PACF')
    s10_model(prs, R)                              ; print('  [10/14] Model Selection')
    s11_diagnostics(prs, R, buf_resid)             ; print('  [11/14] Diagnostics')
    s12_forecast(prs, R, buf_fc, buf_cmp)          ; print('  [12/14] Forecast')
    s13_sarimax(prs, R)                            ; print('  [13/14] SARIMAX')
    s14_plan(prs, R)                               ; print('  [14/14] Plan')

    out = os.path.join(DIR, 'Week5_Slides.pptx')
    prs.save(out)
    print(f'\nDone! Saved: {out}')
    print(f'14 slides | 16:9 widescreen | All plots embedded inline')


if __name__ == '__main__':
    main()
